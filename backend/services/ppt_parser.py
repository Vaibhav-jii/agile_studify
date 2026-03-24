"""
PPT parsing service using python-pptx.
Extracts text content, image counts, and slide-level details from .pptx files.
"""

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import IO
import logging

logger = logging.getLogger(__name__)


class SlideInfo:
    """Parsed information for a single slide."""

    def __init__(self, slide_number: int):
        self.slide_number = slide_number
        self.text_content: str = ""
        self.word_count: int = 0
        self.char_count: int = 0
        self.image_count: int = 0
        self.has_text: bool = False
        self.has_images: bool = False


class PptParseResult:
    """Complete parse result for a PPT file."""

    def __init__(self):
        self.slides: list[SlideInfo] = []
        self.total_slides: int = 0
        self.total_word_count: int = 0
        self.total_char_count: int = 0
        self.total_image_count: int = 0

    @property
    def full_text(self) -> str:
        """Concatenated text from all slides for AI analysis."""
        return "\n\n".join(
            f"[Slide {s.slide_number}] {s.text_content}"
            for s in self.slides
            if s.text_content
        )


def parse_pptx(file_path: str) -> PptParseResult:
    """
    Parse a .pptx file and extract content details.

    Args:
        file_path: Path to the .pptx file on disk.

    Returns:
        PptParseResult with slide-by-slide breakdown.
    """
    result = PptParseResult()

    try:
        prs = Presentation(file_path)
    except Exception as e:
        logger.error(f"Failed to open pptx file: {e}")
        raise ValueError(f"Could not parse the PowerPoint file: {e}")

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_info = SlideInfo(slide_number=slide_index)
        slide_texts: list[str] = []

        for shape in slide.shapes:
            # Extract text from any shape with a text frame
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # Extract text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            slide_texts.append(text)

            # Count images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_info.image_count += 1

            # Count grouped shapes (may contain images)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                slide_info.image_count += _count_images_in_group(shape)

        # Aggregate slide text
        full_text = " ".join(slide_texts)
        slide_info.text_content = full_text
        slide_info.word_count = len(full_text.split()) if full_text else 0
        slide_info.char_count = len(full_text)
        slide_info.has_text = slide_info.word_count > 0
        slide_info.has_images = slide_info.image_count > 0

        result.slides.append(slide_info)

    # Totals
    result.total_slides = len(result.slides)
    result.total_word_count = sum(s.word_count for s in result.slides)
    result.total_char_count = sum(s.char_count for s in result.slides)
    result.total_image_count = sum(s.image_count for s in result.slides)

    return result


def _count_images_in_group(group_shape) -> int:
    """Recursively count images inside grouped shapes."""
    count = 0
    try:
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                count += _count_images_in_group(shape)
    except Exception:
        pass
    return count
